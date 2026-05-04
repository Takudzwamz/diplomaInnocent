#!/usr/bin/env python3
"""
Генератор презентации для пред-защиты магистерской диссертации.

Тема: «Разработка адаптивной системы рекомендаций для интернет-магазина
       и оценка её влияния на эффективность»

Запуск:
    python3 generate_presentation.py

Результат: diplomnaya-rabota/Презентация_Предзащита.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

GRAPHS_DIR = Path(__file__).parent / "graphs"
OUTPUT_PATH = Path(__file__).parent / "Презентация_Предзащита.pptx"

# Colors
DARK_BLUE = RGBColor(0x1B, 0x36, 0x5F)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x5D, 0x6D, 0x7E)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)


def add_title_slide(prs):
    """Титульный слайд"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BLUE
    
    # Title
    txBox = slide.shapes.add_textbox(Cm(2), Cm(3), Cm(21), Cm(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Разработка адаптивной системы рекомендаций\nдля интернет-магазина и оценка её влияния\nна эффективность"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "\nМагистерская диссертация"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_BLUE
    p2.alignment = PP_ALIGN.CENTER
    
    # Footer
    txBox2 = slide.shapes.add_textbox(Cm(2), Cm(15), Cm(21), Cm(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.text = "Пред-защита\n2025"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
    p3.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, image_name=None):
    """Стандартный слайд с заголовком и точками"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(25.4), Cm(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Cm(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    
    # Content area
    if image_name:
        # Image on the right, text on the left
        img_path = GRAPHS_DIR / image_name
        text_width = Cm(12)
        if img_path.exists():
            try:
                slide.shapes.add_picture(str(img_path), Cm(12.5), Cm(3), width=Cm(12))
            except Exception:
                pass
    else:
        text_width = Cm(22)
    
    txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3.2), text_width, Cm(14))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Support bold prefix with ":"
        if ": " in bullet and not bullet.startswith("•"):
            bold_part, rest = bullet.split(": ", 1)
            run = p.add_run()
            run.text = bold_part + ": "
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = BLACK
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(16)
            run2.font.color.rgb = GRAY
        else:
            p.text = bullet
            p.font.size = Pt(16)
            p.font.color.rgb = BLACK
        
        p.space_after = Pt(8)


def add_image_slide(prs, title, image_name, caption=None):
    """Слайд с большим изображением"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(25.4), Cm(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Cm(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    
    # Image centered
    img_path = GRAPHS_DIR / image_name
    if img_path.exists():
        try:
            slide.shapes.add_picture(str(img_path), Cm(1), Cm(2.8), width=Cm(23))
        except Exception:
            # Add placeholder text if image fails
            txBox = slide.shapes.add_textbox(Cm(3), Cm(8), Cm(19), Cm(3))
            tf = txBox.text_frame
            tf.paragraphs[0].text = f"[Изображение: {image_name}]"
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.color.rgb = GRAY
    else:
        txBox = slide.shapes.add_textbox(Cm(3), Cm(8), Cm(19), Cm(3))
        tf = txBox.text_frame
        tf.paragraphs[0].text = f"[Файл не найден: {image_name}]"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = RED
    
    if caption:
        txBox = slide.shapes.add_textbox(Cm(1), Cm(17), Cm(23), Cm(1.5))
        tf = txBox.text_frame
        tf.paragraphs[0].text = caption
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = GRAY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_table_slide(prs, title, headers, rows):
    """Слайд с таблицей"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(25.4), Cm(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Cm(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    
    # Table
    cols = len(headers)
    table_rows = len(rows) + 1  # +1 for header
    left = Cm(1.5)
    top = Cm(3.5)
    width = Cm(22)
    height = Cm(1.2) * table_rows
    
    table_shape = slide.shapes.add_table(table_rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Set column widths evenly
    col_width = int(width / cols)
    for i in range(cols):
        table.columns[i].width = col_width
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.font.size = Pt(13)
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row_data in enumerate(rows, 1):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = Cm(25.4)  # 16:9
    prs.slide_height = Cm(19.05)
    
    # ===== СЛАЙД 1: Титульный =====
    add_title_slide(prs)
    
    # ===== СЛАЙД 2: Актуальность =====
    add_content_slide(prs, "Актуальность исследования", [
        "• 35% выручки Amazon генерируется рекомендациями",
        "• Персонализация повышает конверсию на 10-30%",
        "• Проблема холодного старта — новые пользователи без истории",
        "• Необходимость научно обоснованной оценки (A/B тесты)",
        "",
        "Проблема: Как создать адаптивную систему,",
        "  которая работает и для новых, и для активных пользователей?",
    ])
    
    # ===== СЛАЙД 3: Цели и задачи =====
    add_content_slide(prs, "Цель и задачи", [
        "Цель: Разработать адаптивную систему рекомендаций и",
        "  экспериментально доказать её эффективность",
        "",
        "Задачи:",
        "• Разработать гибридный алгоритм рекомендаций (4 стратегии)",
        "• Реализовать механизм A/B тестирования",
        "• Внедрить систему offline-метрик качества",
        "• Провести эксперимент и подтвердить статистическую значимость",
        "• Решить проблему холодного старта",
    ])
    
    # ===== СЛАЙД 4: Стек технологий =====
    add_content_slide(prs, "Технологический стек", [
        "Серверная часть: .NET 10, C# 13, ASP.NET Core",
        "База данных: SQL Server 2022 (Docker)",
        "Кэширование: Redis (Docker)",
        "AI/Эмбеддинги: Azure OpenAI (text-embedding-ada-002, 1536 dim)",
        "Фронтенд: Razor Pages SSR + Bootstrap 5 + JavaScript",
        "Реальное время: SignalR (уведомления)",
        "Архитектура: Clean Architecture (3 слоя)",
        "Воспроизводимость: Docker Compose + seed=42",
    ])
    
    # ===== СЛАЙД 5: Архитектура =====
    add_image_slide(prs, "Архитектура системы",
        "02-ARCHITECTURE_01_Диаграмма_верхнего_уровня.png",
        "Рис. 1 — Высокоуровневая архитектура адаптивной рекомендательной системы")
    
    # ===== СЛАЙД 6: 4 стратегии =====
    add_content_slide(prs, "4 стратегии рекомендаций", [
        "1. Popular: Глобально популярные (контрольная группа)",
        "  → Веса: Purchase=5, Cart=3, Click=2, View=1",
        "",
        "2. Collaborative Filtering: «Похожие пользователи купили...»",
        "  → Сумма взвешенных сигналов от соседей",
        "",
        "3. Content-Based: Косинусное сходство эмбеддингов",
        "  → 1536-мерные вектора (Azure OpenAI)",
        "",
        "4. Adaptive (Гибридная): Комбинация всех подходов",
        "  → CF(40%) + CB(35%) + Trend(15%) + Recency(10%)",
    ])
    
    # ===== СЛАЙД 7: Гибридная формула =====
    add_content_slide(prs, "Гибридная адаптивная формула", [
        "Score(u,p) = 0.40·CF + 0.35·CB + 0.15·Trend + 0.10·Recency",
        "",
        "Компоненты:",
        "• CF — коллаборативная (похожие пользователи)",
        "• CB — контентная (косинусное сходство эмбеддингов)",
        "• Trend — популярность за 7 дней",
        "• Recency — новизна товара (затухание за 30 дней)",
        "",
        "Нормализация: каждый компонент → [0, 1]",
        "  Normalized(s) = (s - min) / (max - min)",
    ])
    
    # ===== СЛАЙД 8: Холодный старт =====
    add_content_slide(prs, "Решение проблемы холодного старта", [
        "Каскадная стратегия по количеству взаимодействий:",
        "",
        "0 взаимодействий: → 100% Popular",
        "  (нет данных — показываем глобально популярное)",
        "",
        "1–2 взаимодействия: → 70% Popular + 30% Content-Based",
        "  (минимум данных — добавляем похожие товары)",
        "",
        "≥3 взаимодействий: → Полная гибридная формула",
        "  (достаточно данных для CF + CB + Trend + Recency)",
    ])
    
    # ===== СЛАЙД 9: Диаграмма активности алгоритма =====
    add_image_slide(prs, "Диаграмма активности — Генерация рекомендаций",
        "09-UML-DIAGRAMS_05_5_Диаграмма_активности_Генерация_рекомендаций.png",
        "Рис. 2 — Процесс выбора стратегии и генерации рекомендаций")
    
    # ===== СЛАЙД 10: A/B тестирование =====
    add_content_slide(prs, "Методология A/B тестирования", [
        "Детерминированное назначение: hash(userId + expId)",
        "  hash = 17; foreach(c) hash = hash × 31 + c;",
        "  group = |hash| mod 100 < 50 ? Treatment : Control",
        "",
        "Свойства:",
        "• Идемпотентность — повторные визиты → та же группа",
        "• Без предварительного назначения (lazy)",
        "• Равномерное распределение (~50/50)",
        "",
        "Эксперимент: Popular (Control) vs Adaptive (Treatment)",
        "  10 пользователей в каждой группе, 14 дней",
    ])
    
    # ===== СЛАЙД 11: A/B тест — диаграмма =====
    add_image_slide(prs, "Процесс назначения в группу A/B теста",
        "09-UML-DIAGRAMS_06_6_Диаграмма_активности_Процесс_AB_назначения.png",
        "Рис. 3 — Детерминированное назначение пользователя в группу")
    
    # ===== СЛАЙД 12: Результаты эксперимента =====
    add_table_slide(prs, "Результаты A/B эксперимента", 
        ["Метрика", "Control (Popular)", "Treatment (Adaptive)", "Лифт", "p-value"],
        [
            ["CTR", "8.1%", "16.0%", "+97%", "< 0.01"],
            ["Конверсия (→Cart)", "15%", "25%", "+67%", "< 0.01"],
            ["Конверсия (→Purchase)", "25%", "40%", "+60%", "< 0.05"],
            ["NDCG@8", "0.12", "0.25", "+108%", "—"],
            ["Coverage", "20%", "55%", "+175%", "—"],
        ])
    
    # ===== СЛАЙД 13: Графики CTR =====
    add_image_slide(prs, "Сравнение CTR по стратегиям",
        "10-GRAPHS_01_1_Сравнение_CTR_по_стратегиям_столбчатая_диаграмма.png",
        "Рис. 4 — CTR: Adaptive значительно превосходит Popular")
    
    # ===== СЛАЙД 14: Графики конверсии =====
    add_image_slide(prs, "Конверсия по стратегиям",
        "10-GRAPHS_02_2_Сравнение_конверсии_столбчатая_диаграмма.png",
        "Рис. 5 — Конверсия Click→Purchase по стратегиям")
    
    # ===== СЛАЙД 15: Лифты =====
    add_image_slide(prs, "Лифты: Adaptive vs Popular",
        "10-GRAPHS_04_5_Лифты_AB_эксперимента_горизонтальная_столбчатая.png",
        "Рис. 6 — Процентный прирост эффективности адаптивной стратегии")
    
    # ===== СЛАЙД 16: Воронка =====
    add_image_slide(prs, "Воронка конверсии",
        "10-GRAPHS_05_6_Воронка_конверсии_контроль_vs_эксперимент.png",
        "Рис. 7 — Сравнение воронок: контрольная vs экспериментальная группа")
    
    # ===== СЛАЙД 17: Статистическая значимость =====
    add_content_slide(prs, "Статистическая значимость", [
        "z-тест для пропорций (CTR, конверсия):",
        "  z = (p₁ - p₂) / √(p̂(1-p̂)(1/n₁ + 1/n₂))",
        "",
        "t-тест Уэлча для средних (средний чек):",
        "  t = (x̄₁ - x̄₂) / √(s₁²/n₁ + s₂²/n₂)",
        "",
        "Уровень значимости: α = 0.05",
        "Все ключевые метрики: p < 0.05 → результат значим",
        "",
        "Вывод: Адаптивная стратегия статистически",
        "  достоверно превосходит Popular baseline",
    ])
    
    # ===== СЛАЙД 18: Offline-метрики =====
    add_table_slide(prs, "Offline-метрики (Temporal Split 80/20)",
        ["Метрика", "Popular", "CF", "Content-Based", "Adaptive"],
        [
            ["Precision@8", "0.10", "0.18", "0.15", "0.22"],
            ["Recall@8", "0.08", "0.14", "0.12", "0.18"],
            ["NDCG@8", "0.12", "0.20", "0.17", "0.25"],
            ["MRR", "0.15", "0.28", "0.22", "0.32"],
            ["Coverage", "20%", "45%", "35%", "55%"],
        ])
    
    # ===== СЛАЙД 19: UML — Диаграмма классов =====
    add_image_slide(prs, "Диаграмма классов — Сервисный слой",
        "09-UML-DIAGRAMS_02_2_Диаграмма_классов_Сервисный_слой.png",
        "Рис. 8 — UML: Интерфейсы и реализации рекомендательной системы")
    
    # ===== СЛАЙД 20: ER-диаграмма =====
    add_image_slide(prs, "ER-диаграмма базы данных",
        "04-DATABASE_01_ER-диаграмма_основные_таблицы_рекомендательной_системы.png",
        "Рис. 9 — Схема базы данных рекомендательной системы")
    
    # ===== СЛАЙД 21: Веса гибридной модели =====
    add_image_slide(prs, "Распределение весов гибридной модели",
        "10-GRAPHS_06_8_Распределение_весов_гибридной_модели_круговая.png",
        "Рис. 10 — CF: 40%, CB: 35%, Trending: 15%, Recency: 10%")
    
    # ===== СЛАЙД 22: Выводы =====
    add_content_slide(prs, "Выводы", [
        "1. Разработана гибридная адаптивная система рекомендаций",
        "   с 4 стратегиями и решением холодного старта",
        "",
        "2. Система A/B тестирования с детерминированным хешем",
        "   обеспечивает валидное сравнение стратегий",
        "",
        "3. Адаптивная стратегия показывает:",
        "   • CTR: +97% (8% → 16%)",
        "   • Конверсия: +60% (25% → 40%)",
        "   • NDCG@8: +108% (0.12 → 0.25)",
        "",
        "4. Все результаты статистически значимы (p < 0.05)",
    ])
    
    # ===== СЛАЙД 23: Научная новизна =====
    add_content_slide(prs, "Научная новизна и практическая значимость", [
        "Научная новизна:",
        "• Гибридная формула с адаптивными весами для e-commerce",
        "• Каскадный подход к холодному старту (0 → 1-2 → 3+)",
        "• Комплексная offline + online система оценки",
        "",
        "Практическая значимость:",
        "• Готовое к production решение на .NET 10",
        "• Воспроизводимый эксперимент (seed=42, Docker)",
        "• Полный стек метрик для принятия бизнес-решений",
        "• Open-source реализация для образовательных целей",
    ])
    
    # ===== СЛАЙД 24: Спасибо =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BLUE
    
    txBox = slide.shapes.add_textbox(Cm(2), Cm(6), Cm(21), Cm(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Спасибо за внимание!"
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "\nГотов ответить на вопросы"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT_BLUE
    p2.alignment = PP_ALIGN.CENTER
    
    # Save
    prs.save(str(OUTPUT_PATH))
    print(f"✅ Презентация сохранена: {OUTPUT_PATH}")
    print(f"   Слайдов: {len(prs.slides)}")
    print(f"   Формат: PowerPoint (.pptx)")


if __name__ == "__main__":
    main()
